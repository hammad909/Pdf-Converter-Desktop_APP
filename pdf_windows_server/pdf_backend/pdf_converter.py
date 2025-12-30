from fastapi import APIRouter, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import pdfplumber
import uuid
import io
from pathlib import Path
from pypdf import PdfReader, PdfWriter 
from paths import UPLOAD_DIR, CONVERTED_DIR

router = APIRouter()

PPTX_DIR = CONVERTED_DIR


def group_words_into_lines(words):
    lines = []
    current_line = []
    prev_y = None
    for word in words:
        y = round(word["top"], 1)
        if prev_y is not None and abs(y - prev_y) > 2:
            lines.append(current_line)
            current_line = []
        current_line.append(word)
        prev_y = y
    if current_line:
        lines.append(current_line)
    return lines


@router.post("/upload/ppt/")
async def convert_pdf_to_pptx(file: UploadFile = File(...)):
    if not file.filename.endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are supported")

    file_id = str(uuid.uuid4())
    pdf_path = UPLOAD_DIR / f"{file_id}.pdf"
    pptx_path = PPTX_DIR / f"{file_id}.pptx"
    decrypted_pdf_path = UPLOAD_DIR / f"{file_id}_decrypted.pdf"

    # Save uploaded file
    with open(pdf_path, "wb") as f:
        f.write(await file.read())

    try:
        # ✅ Step 1: Try to decrypt the PDF if needed
        try:
            reader = PdfReader(str(pdf_path))
            if reader.is_encrypted:
                try:
                    reader.decrypt("")  # Try empty password
                    writer = PdfWriter()
                    for page in reader.pages:
                        writer.add_page(page)
                    with open(decrypted_pdf_path, "wb") as out_file:
                        writer.write(out_file)
                    pdf_to_open = decrypted_pdf_path
                except Exception:
                    raise HTTPException(status_code=400, detail="PDF is encrypted and cannot be decrypted.")
            else:
                pdf_to_open = pdf_path
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error reading PDF: {e}")

        # ✅ Step 2: Normal conversion (your logic)
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        MAX_HEIGHT = Inches(6.5)
        LEFT = Inches(1)
        TOP_START = Inches(1)
        WIDTH = Inches(8)

        with pdfplumber.open(str(pdf_to_open)) as pdf:
            for page_index, page in enumerate(pdf.pages, start=1):
                slide = prs.slides.add_slide(blank_slide_layout)
                top = TOP_START
                content_blocks = []

                # Extract text
                words = page.extract_words(extra_attrs=["fontname", "size"])
                lines = group_words_into_lines(words) if words else []
                for line in lines:
                    font_size = float(line[0].get("size", 12))
                    height = Inches(font_size * 1.5 / 72)
                    content_blocks.append({
                        "type": "text",
                        "line": line,
                        "height": height
                    })

                # Extract images
                for img in page.images:
                    try:
                        x0, top_img, x1, bottom = img["x0"], img["top"], img["x1"], img["bottom"]
                        bbox = (x0, top_img, x1, bottom)
                        cropped_image = page.crop(bbox).to_image(resolution=150).original

                        img_byte_arr = io.BytesIO()
                        cropped_image.save(img_byte_arr, format='PNG')
                        img_byte_arr.seek(0)

                        pil_img = Image.open(img_byte_arr)
                        aspect_ratio = pil_img.height / pil_img.width
                        img_width = Inches(5)
                        img_height = img_width * aspect_ratio

                        content_blocks.append({
                            "type": "image",
                            "data": img_byte_arr,
                            "height": img_height
                        })
                    except Exception as e:
                        print(f"Image extract failed on page {page_index}:", e)

                # ✅ NEW: handle image-only pages (no text, no extractable images)
                if not lines and not page.images:
                    try:
                        full_img = page.to_image(resolution=150).original
                        img_byte_arr = io.BytesIO()
                        full_img.save(img_byte_arr, format='PNG')
                        img_byte_arr.seek(0)

                        pil_img = Image.open(img_byte_arr)
                        aspect_ratio = pil_img.height / pil_img.width
                        img_width = Inches(7)
                        img_height = img_width * aspect_ratio

                        content_blocks.append({
                            "type": "image",
                            "data": img_byte_arr,
                            "height": img_height
                        })
                    except Exception as e:
                        print(f"Full-page image extract failed on page {page_index}:", e)

                # Render blocks
                for block in content_blocks:
                    height = block["height"]
                    if top + height > MAX_HEIGHT:
                        slide = prs.slides.add_slide(blank_slide_layout)
                        top = TOP_START

                    if block["type"] == "text":
                        line = block["line"]
                        textbox = slide.shapes.add_textbox(LEFT, top, WIDTH, height).text_frame
                        p = textbox.paragraphs[0]
                        for word in line:
                            run = p.add_run()
                            run.text = word['text'] + " "
                            run.font.size = Pt(float(word["size"]))
                            fontname = word.get("fontname", "").lower()
                            run.font.bold = "bold" in fontname
                            run.font.italic = "italic" in fontname or "oblique" in fontname
                            run.font.color.rgb = RGBColor(0, 0, 0)
                        top += height + Inches(0.05)

                    elif block["type"] == "image":
                        try:
                            if block["height"] > MAX_HEIGHT - TOP_START:
                                slide = prs.slides.add_slide(blank_slide_layout)
                                top = TOP_START
                            slide.shapes.add_picture(block["data"], LEFT, top, width=Inches(5), height=block["height"])
                            top += block["height"] + Inches(0.2)
                        except Exception as e:
                            print(f"Image insert failed on page {page_index}:", e)

        prs.save(str(pptx_path))

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PPTX conversion failed: {e}")

    return {
        "message": "PDF converted to PPTX successfully",
        "file": pptx_path.name,
        "download_url": f"/download/{pptx_path.name}"
    }
